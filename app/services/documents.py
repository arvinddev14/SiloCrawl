"""Document intelligence: PDFs, Office files, CSV, and images -> clean text.

One converter per format, all OSS and pip-installable. OCR (images and the
scanned-PDF case) is an optional extra — ``pip install silocrawl[ocr]`` pulls
RapidOCR; without it, image inputs get a clear error instead of a crash.
Downloads honor the same robots.txt + per-domain politeness as the scraper.
"""
from __future__ import annotations

import csv as csv_module
import io
import logging
import zipfile
from pathlib import PurePosixPath
from typing import Any, Callable
from urllib.parse import urlparse

import httpx

from app.core.config import get_settings
from app.models.schemas import OutputFormat
from app.services import cleaner, netguard, robots
from app.services.throttle import throttle

logger = logging.getLogger("silocrawl")
settings = get_settings()

MAX_TABLE_ROWS = 500  # cap per sheet/CSV so a huge workbook can't flood the LLM


class DocumentError(ValueError):
    """User-facing document problem (route maps it to 422)."""


class DocumentTooLargeError(DocumentError):
    """Document exceeds ``document_max_bytes`` (route maps it to 413)."""


# ---------- format detection ----------

_CONTENT_TYPES = {
    "application/pdf": "pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx",
    "text/csv": "csv",
    "application/csv": "csv",
    "text/plain": "txt",
    "text/markdown": "md",
    "text/html": "html",
    "image/png": "image",
    "image/jpeg": "image",
    "image/webp": "image",
}

_EXTENSIONS = {
    ".pdf": "pdf",
    ".docx": "docx",
    ".pptx": "pptx",
    ".xlsx": "xlsx",
    ".csv": "csv",
    ".txt": "txt",
    ".md": "md",
    ".markdown": "md",
    ".html": "html",
    ".htm": "html",
    ".png": "image",
    ".jpg": "image",
    ".jpeg": "image",
    ".webp": "image",
}


def _detect_ooxml(data: bytes) -> str | None:
    """Disambiguate the OOXML zip family by its internal folder layout."""
    try:
        names = zipfile.ZipFile(io.BytesIO(data)).namelist()
    except zipfile.BadZipFile:
        return None
    for prefix, fmt in (("word/", "docx"), ("ppt/", "pptx"), ("xl/", "xlsx")):
        if any(n.startswith(prefix) for n in names):
            return fmt
    return None


def detect_format(
    data: bytes, content_type: str | None = None, filename: str | None = None
) -> str:
    if content_type:
        ct = content_type.split(";")[0].strip().lower()
        if ct in _CONTENT_TYPES:
            return _CONTENT_TYPES[ct]
    if filename:
        ext = PurePosixPath(filename.lower()).suffix
        if ext in _EXTENSIONS:
            return _EXTENSIONS[ext]
    # magic bytes
    if data.startswith(b"%PDF"):
        return "pdf"
    if data.startswith((b"\x89PNG", b"\xff\xd8", b"RIFF")):
        return "image"
    if data.startswith(b"PK\x03\x04"):
        fmt = _detect_ooxml(data)
        if fmt:
            return fmt
    # Last resort: treat as plain text, but only if it plausibly is text
    # (NUL bytes decode fine as UTF-8, so check for them explicitly).
    if b"\x00" not in data[:4096]:
        try:
            data.decode("utf-8")
            return "txt"
        except UnicodeDecodeError:
            pass
    raise DocumentError(
        "Unsupported document type. Supported: pdf, docx, pptx, xlsx, csv, "
        "txt, md, html, png, jpg, webp."
    )


# ---------- OCR (optional extra) ----------

_ocr_engine: Any = None
_ocr_checked = False


def ocr_image(data: bytes) -> str | None:
    """Recognized text, '' if none found, or None when the extra isn't installed."""
    global _ocr_engine, _ocr_checked
    if not _ocr_checked:
        _ocr_checked = True
        try:
            from rapidocr_onnxruntime import RapidOCR

            _ocr_engine = RapidOCR()
        except ImportError:
            _ocr_engine = None
    if _ocr_engine is None:
        return None
    result, _elapsed = _ocr_engine(data)
    if not result:
        return ""
    return "\n".join(item[1] for item in result)


# ---------- converters ----------

def _rows_to_markdown(rows: list[list[str]]) -> str:
    if not rows:
        return ""
    header, *body = rows
    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join("---" for _ in header) + " |",
    ]
    lines += ["| " + " | ".join(row) + " |" for row in body]
    return "\n".join(lines)


def _pdf_convert(data: bytes) -> tuple[str, dict[str, Any]]:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    parts: list[str] = []
    ocr_pages = 0
    for page in reader.pages:
        text = (page.extract_text() or "").strip()
        if not text:  # scanned page — try OCR on its embedded images
            recovered: list[str] = []
            try:
                for image in page.images:
                    ocr_text = ocr_image(image.data)
                    if ocr_text:
                        recovered.append(ocr_text)
            except Exception:  # noqa: BLE001 - broken embedded images happen
                logger.warning("pdf_image_ocr_failed", exc_info=True)
            if recovered:
                text = "\n".join(recovered)
                ocr_pages += 1
        if text:
            parts.append(text)
    if not parts:
        raise DocumentError(
            "No extractable text found. If this is a scanned PDF, install the "
            "OCR extra: pip install silocrawl[ocr]"
        )
    return "\n\n".join(parts), {"pages": len(reader.pages), "ocr_pages": ocr_pages}


def _docx_convert(data: bytes) -> tuple[str, dict[str, Any]]:
    from docx import Document

    doc = Document(io.BytesIO(data))
    parts = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    paragraphs = len(parts)
    for table in doc.tables:
        parts.append(
            _rows_to_markdown(
                [[cell.text.strip() for cell in row.cells] for row in table.rows]
            )
        )
    return "\n\n".join(parts), {"paragraphs": paragraphs, "tables": len(doc.tables)}


def _pptx_convert(data: bytes) -> tuple[str, dict[str, Any]]:
    from pptx import Presentation

    pres = Presentation(io.BytesIO(data))
    slides: list[str] = []
    for i, slide in enumerate(pres.slides, 1):
        texts = [
            shape.text_frame.text.strip()
            for shape in slide.shapes
            if shape.has_text_frame and shape.text_frame.text.strip()
        ]
        slides.append(f"## Slide {i}\n\n" + "\n\n".join(texts))
    return "\n\n".join(slides), {"slides": len(slides)}


def _xlsx_convert(data: bytes) -> tuple[str, dict[str, Any]]:
    from openpyxl import load_workbook

    workbook = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts: list[str] = []
    for sheet in workbook.worksheets:
        rows: list[list[str]] = []
        for row in sheet.iter_rows(values_only=True):
            rows.append(["" if v is None else str(v) for v in row])
            if len(rows) >= MAX_TABLE_ROWS:
                break
        if rows:
            parts.append(f"## {sheet.title}\n\n" + _rows_to_markdown(rows))
    return "\n\n".join(parts), {"sheets": len(workbook.worksheets)}


def _csv_convert(data: bytes) -> tuple[str, dict[str, Any]]:
    text = data.decode("utf-8", errors="replace")
    rows = [list(r) for r in csv_module.reader(io.StringIO(text))][:MAX_TABLE_ROWS]
    return _rows_to_markdown(rows), {"rows": len(rows)}


def _text_convert(data: bytes) -> tuple[str, dict[str, Any]]:
    return data.decode("utf-8", errors="replace"), {}


def _html_convert(data: bytes) -> tuple[str, dict[str, Any]]:
    result = cleaner.clean(
        html=data.decode("utf-8", errors="replace"),
        url="",
        status_code=200,
        formats=[OutputFormat.markdown],
        only_main_content=True,
        include_tags=None,
        exclude_tags=None,
    )
    return result.markdown or "", {}


def _image_convert(data: bytes) -> tuple[str, dict[str, Any]]:
    text = ocr_image(data)
    if text is None:
        raise DocumentError(
            "Image input requires the OCR extra: pip install silocrawl[ocr]"
        )
    return text, {"ocr": True}


_CONVERTERS: dict[str, Callable[[bytes], tuple[str, dict[str, Any]]]] = {
    "pdf": _pdf_convert,
    "docx": _docx_convert,
    "pptx": _pptx_convert,
    "xlsx": _xlsx_convert,
    "csv": _csv_convert,
    "txt": _text_convert,
    "md": _text_convert,
    "html": _html_convert,
    "image": _image_convert,
}


def convert(
    data: bytes, *, content_type: str | None = None, filename: str | None = None
) -> tuple[str, dict[str, Any]]:
    fmt = detect_format(data, content_type, filename)
    text, meta = _CONVERTERS[fmt](data)
    return text, {"format": fmt, "size_bytes": len(data), **meta}


# ---------- entry points ----------

async def download(url: str) -> tuple[bytes, str | None]:
    """Fetch document bytes with the same politeness rules as the scraper."""
    min_delay: float | None = None
    if settings.respect_robots:
        await robots.check(url)
        site_delay = await robots.crawl_delay(url)
        if site_delay is not None:
            min_delay = max(site_delay, settings.per_domain_delay)
    await throttle.wait(url, min_delay)

    async with httpx.AsyncClient(
        follow_redirects=True,
        timeout=settings.request_timeout,
        headers={"User-Agent": settings.user_agent},
        event_hooks=netguard.event_hooks(),
    ) as client:
        # Stream with a cap: the size check must happen while downloading, not
        # after the whole body already sits in memory.
        async with client.stream("GET", url) as resp:
            resp.raise_for_status()
            limit = settings.document_max_bytes
            declared = resp.headers.get("content-length", "")
            if declared.isdigit() and int(declared) > limit:
                raise DocumentTooLargeError(
                    f"Document declares {declared} bytes; limit is {limit}."
                )
            chunks: list[bytes] = []
            total = 0
            async for chunk in resp.aiter_bytes():
                total += len(chunk)
                if total > limit:
                    raise DocumentTooLargeError(
                        f"Document exceeded the {limit}-byte limit."
                    )
                chunks.append(chunk)
        return b"".join(chunks), resp.headers.get("content-type")


async def process(
    *,
    url: str | None = None,
    data: bytes | None = None,
    content_type: str | None = None,
    filename: str | None = None,
) -> tuple[str, dict[str, Any]]:
    """Resolve (download if needed), size-check, and convert a document."""
    if url:
        data, content_type = await download(url)
        filename = filename or urlparse(url).path
    if data is None:
        raise DocumentError("Provide a document URL or file.")
    if len(data) > settings.document_max_bytes:
        raise DocumentTooLargeError(
            f"Document is {len(data)} bytes; limit is {settings.document_max_bytes}."
        )
    text, meta = convert(data, content_type=content_type, filename=filename)
    meta["source"] = url or filename or "upload"
    return text, meta
