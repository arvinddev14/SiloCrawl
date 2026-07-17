import io
import json

import pytest
import respx
from httpx import Response

from app.llm import get_router
from app.llm.base import LLMResponse, ToolCall
from app.services import documents


# ---------- fixtures / builders ----------

def _make_pdf(text: str) -> bytes:
    """Assemble a minimal but valid one-page PDF containing ``text``."""
    stream = f"BT /F1 24 Tf 72 720 Td ({text}) Tj ET".encode()
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Contents 4 0 R "
        b"/Resources << /Font << /F1 5 0 R >> >> >>",
        b"<< /Length " + str(len(stream)).encode() + b" >>\nstream\n" + stream + b"\nendstream",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
    ]
    out = bytearray(b"%PDF-1.4\n")
    offsets = []
    for i, obj in enumerate(objects, 1):
        offsets.append(len(out))
        out += f"{i} 0 obj\n".encode() + obj + b"\nendobj\n"
    xref = len(out)
    out += f"xref\n0 {len(objects) + 1}\n".encode() + b"0000000000 65535 f \n"
    for off in offsets:
        out += f"{off:010d} 00000 n \n".encode()
    out += (
        f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\n"
        f"startxref\n{xref}\n%%EOF"
    ).encode()
    return bytes(out)


def _make_docx() -> bytes:
    from docx import Document

    doc = Document()
    doc.add_paragraph("Quarterly report intro.")
    table = doc.add_table(rows=2, cols=2)
    table.rows[0].cells[0].text = "Region"
    table.rows[0].cells[1].text = "Revenue"
    table.rows[1].cells[0].text = "EMEA"
    table.rows[1].cells[1].text = "42"
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _make_xlsx() -> bytes:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Sales"
    ws.append(["Region", "Revenue"])
    ws.append(["EMEA", 42])
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _make_pptx() -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    pres = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[6])  # blank layout
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "SiloLoop roadmap"
    buf = io.BytesIO()
    pres.save(buf)
    return buf.getvalue()


# ---------- detection ----------

def test_detect_by_content_type():
    assert documents.detect_format(b"", "application/pdf; charset=x") == "pdf"
    assert documents.detect_format(b"", "text/csv") == "csv"


def test_detect_by_extension():
    assert documents.detect_format(b"", None, "report.XLSX".lower()) == "xlsx"
    assert documents.detect_format(b"", None, "notes.md") == "md"


def test_detect_by_magic_bytes():
    assert documents.detect_format(b"%PDF-1.7 rest") == "pdf"
    assert documents.detect_format(b"\x89PNG\r\n\x1a\n...") == "image"
    assert documents.detect_format(_make_docx()) == "docx"  # OOXML zip sniffing
    assert documents.detect_format(_make_xlsx()) == "xlsx"


def test_detect_rejects_unknown_binary():
    with pytest.raises(documents.DocumentError):
        documents.detect_format(b"\x00\x01\x02\x03binary")


# ---------- converters ----------

def test_csv_to_markdown_table():
    text, meta = documents.convert(b"a,b\n1,2\n3,4", content_type="text/csv")
    assert "| a | b |" in text
    assert "| 1 | 2 |" in text
    assert meta["format"] == "csv"
    assert meta["rows"] == 3


def test_csv_row_cap(monkeypatch):
    monkeypatch.setattr(documents, "MAX_TABLE_ROWS", 2)
    text, meta = documents.convert(b"a,b\n1,2\n3,4\n5,6", content_type="text/csv")
    assert meta["rows"] == 2
    assert "5" not in text


def test_docx_paragraphs_and_tables():
    text, meta = documents.convert(_make_docx())
    assert "Quarterly report intro." in text
    assert "| Region | Revenue |" in text
    assert meta["tables"] == 1


def test_xlsx_sheet_table():
    text, meta = documents.convert(_make_xlsx())
    assert "## Sales" in text
    assert "| EMEA | 42 |" in text
    assert meta["sheets"] == 1


def test_pptx_slide_text():
    text, meta = documents.convert(_make_pptx())
    assert "## Slide 1" in text
    assert "SiloLoop roadmap" in text
    assert meta["slides"] == 1


def test_pdf_text_extraction():
    text, meta = documents.convert(_make_pdf("Hello World"))
    assert "Hello World" in text
    assert meta["pages"] == 1


def test_html_via_cleaner():
    html = b"<html><body><article><h1>Title</h1><p>Body text.</p></article></body></html>"
    text, meta = documents.convert(html, content_type="text/html")
    assert "Title" in text
    assert meta["format"] == "html"


# ---------- OCR paths ----------

def test_image_without_ocr_extra_errors(monkeypatch):
    monkeypatch.setattr(documents, "ocr_image", lambda data: None)
    with pytest.raises(documents.DocumentError, match="OCR extra"):
        documents.convert(b"\x89PNG\r\n\x1a\nfake")


def test_image_with_ocr(monkeypatch):
    monkeypatch.setattr(documents, "ocr_image", lambda data: "RECOGNIZED TEXT")
    text, meta = documents.convert(b"\x89PNG\r\n\x1a\nfake")
    assert text == "RECOGNIZED TEXT"
    assert meta["ocr"] is True


def test_empty_pdf_without_ocr_gives_clear_error(monkeypatch):
    from pypdf import PdfWriter

    monkeypatch.setattr(documents, "ocr_image", lambda data: None)
    writer = PdfWriter()
    writer.add_blank_page(width=612, height=792)
    buf = io.BytesIO()
    writer.write(buf)
    with pytest.raises(documents.DocumentError, match="No extractable text"):
        documents.convert(buf.getvalue())


# ---------- routes ----------

@pytest.fixture
def no_politeness(monkeypatch):
    monkeypatch.setattr(documents.settings, "respect_robots", False)
    monkeypatch.setattr(documents.settings, "per_domain_delay", 0.0)


@respx.mock
async def test_document_url_route(client, no_politeness):
    respx.get("https://docs.test/report.pdf").mock(
        return_value=Response(
            200, content=_make_pdf("Annual figures"),
            headers={"content-type": "application/pdf"},
        )
    )
    resp = await client.post("/v1/document", json={"url": "https://docs.test/report.pdf"})
    assert resp.status_code == 200
    body = resp.json()
    assert "Annual figures" in body["text"]
    assert body["metadata"]["format"] == "pdf"
    assert body["data"] is None  # no schema -> conversion only


async def test_document_upload_route(client):
    resp = await client.post(
        "/v1/document/upload",
        files={"file": ("data.csv", b"name,qty\nWidget,7", "text/csv")},
    )
    assert resp.status_code == 200
    assert "| Widget | 7 |" in resp.json()["text"]


async def test_document_upload_with_extraction(client, monkeypatch):
    class FakeProvider:
        async def complete(self, **kw):
            return LLMResponse(
                tool_calls=[
                    ToolCall(
                        name="emit_extracted_data",
                        arguments=json.dumps({"name": "Widget", "qty": 7}),
                    )
                ]
            )

    monkeypatch.setattr(get_router()._registry, "provider_for", lambda spec: FakeProvider())
    schema = json.dumps(
        {"type": "object", "properties": {"name": {"type": "string"}, "qty": {"type": "integer"}}}
    )
    resp = await client.post(
        "/v1/document/upload",
        files={"file": ("data.csv", b"name,qty\nWidget,7", "text/csv")},
        data={"schema": schema},
    )
    assert resp.status_code == 200
    assert resp.json()["data"] == {"name": "Widget", "qty": 7}


async def test_document_upload_bad_schema_json(client):
    resp = await client.post(
        "/v1/document/upload",
        files={"file": ("data.csv", b"a,b", "text/csv")},
        data={"schema": "{not json"},
    )
    assert resp.status_code == 400


async def test_document_upload_too_large(client, monkeypatch):
    monkeypatch.setattr(documents.settings, "document_max_bytes", 10)
    resp = await client.post(
        "/v1/document/upload",
        files={"file": ("data.csv", b"a,b\n" * 100, "text/csv")},
    )
    assert resp.status_code == 413


async def test_document_upload_unsupported_type(client):
    resp = await client.post(
        "/v1/document/upload",
        files={"file": ("blob.bin", b"\x00\x01\x02\x03", "application/octet-stream")},
    )
    assert resp.status_code == 422
